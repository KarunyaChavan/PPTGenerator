from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from datetime import datetime
from typing import Dict, List, Tuple
from flask import current_app

class PPTGeneratorService:
    """Service for generating PowerPoint presentations"""

    def __init__(self):
        self.company_colors = {
            'primary': RGBColor(0, 51, 102),
            'secondary': RGBColor(0, 123, 191),
            'accent': RGBColor(255, 127, 0),
            'text': RGBColor(64, 64, 64),
            'light': RGBColor(245, 245, 245)
        }

    def generate_presentation(self, presentation_obj, slides_data: List[Dict]) -> Tuple[str, str]:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        self._create_title_slide(prs, presentation_obj)

        # Agenda
        if presentation_obj.agenda:
            # Support both JSON list or string
            try:
                agenda_items = json.loads(presentation_obj.agenda)
                if isinstance(agenda_items, list):
                    agenda_text = "\n".join(agenda_items)
                else:
                    agenda_text = str(agenda_items)
            except:
                agenda_text = str(presentation_obj.agenda)
            self._create_agenda_slide(prs, agenda_text)

        # Content slides
        if slides_data:
            for slide_data in slides_data:
                self._create_content_slide(prs, slide_data)

        self._create_thank_you_slide(prs)

        filename = self._generate_filename(presentation_obj)
        file_path = self._get_file_path(presentation_obj.id, filename)

        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        prs.save(file_path)
        abs_file_path = os.path.abspath(file_path)  # Ensure absolute path is stored
        return abs_file_path, filename

    def _create_title_slide(self, prs, presentation_obj):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.company_colors['light']

        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = presentation_obj.title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.company_colors['primary']
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle_text = f"Presented by: {presentation_obj.author.username}"
        if presentation_obj.author.department:
            subtitle_text += f" | {presentation_obj.author.department}"
        subtitle_text += f"\nDate: {datetime.now().strftime('%B %d, %Y')}"

        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_frame.paragraphs[0].font.size = Pt(18)
        subtitle_frame.paragraphs[0].font.color.rgb = self.company_colors['text']
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _create_agenda_slide(self, prs, agenda_text: str):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Agenda"
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = self.company_colors['primary']

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        agenda_items = [line.strip() for line in agenda_text.split("\n") if line.strip()]
        for i, item in enumerate(agenda_items, 1):
            p = text_frame.paragraphs[0] if i == 1 else text_frame.add_paragraph()
            p.text = f"{i}. {item}"
            p.font.size = Pt(24)
            p.font.color.rgb = self.company_colors['text']
            p.space_after = Pt(12)

    def _create_content_slide(self, prs, slide_data: Dict):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.get('title', 'Slide Title')
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = self.company_colors['primary']

        content_text = slide_data.get('content', '')
        bullet_points = slide_data.get('bullet_points', [])

        if content_text or bullet_points:
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()

            lines = []
            if content_text:
                lines.extend([line.strip() for line in content_text.split("\n") if line.strip()])
            if bullet_points:
                lines.extend(bullet_points)

            for i, line in enumerate(lines):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = line
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = self.company_colors['text']
                p.space_after = Pt(6)

    def _create_thank_you_slide(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.company_colors['primary']

        thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
        thank_you_frame = thank_you_box.text_frame
        thank_you_frame.text = "Thank You"
        thank_you_frame.paragraphs[0].font.size = Pt(48)
        thank_you_frame.paragraphs[0].font.bold = True
        thank_you_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        thank_you_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _generate_filename(self, presentation_obj) -> str:
        safe_title = "".join(c for c in presentation_obj.title if c.isalnum() or c in (' ', '-', '_')).rstrip()
        safe_title = safe_title.replace(' ', '_')[:50]
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        return f"{safe_title}_{timestamp}.pptx"

    def _get_file_path(self, presentation_id: int, filename: str) -> str:
        storage_dir = current_app.config['UPLOAD_FOLDER']
        presentation_dir = os.path.join(storage_dir, str(presentation_id))
        return os.path.join(presentation_dir, filename)
